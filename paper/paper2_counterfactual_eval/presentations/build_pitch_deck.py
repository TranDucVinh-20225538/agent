#!/usr/bin/env python3
"""Build PPTX for anh Khải pitch: Paper 1 arc + Paper 2 decision experiment."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).resolve().parent / "khai_paper2_pitch.pptx"

TITLE_COLOR = RGBColor(0x1A, 0x36, 0x5D)
ACCENT = RGBColor(0x2E, 0x6B, 0x9E)
BODY = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x55, 0x55, 0x55)


def set_title(shape, text: str, size: int = 28):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.LEFT


def add_bullets(text_frame, lines: list[str], size: int = 18, level0=True):
    text_frame.clear()
    for i, line in enumerate(lines):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = BODY if level0 else MUTED
        p.level = 0
        p.space_after = Pt(8)


def add_slide_title_only(prs, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(8.8), Inches(1.2))
    set_title(box, title, 32)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(3.4), Inches(8.8), Inches(0.8))
        tf = sub.text_frame
        tf.text = subtitle
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = MUTED


def add_content_slide(prs, title: str, bullets: list[str], note: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.9))
    set_title(tbox, title, 24)
    body = slide.shapes.add_textbox(Inches(0.55), Inches(1.25), Inches(8.9), Inches(5.8))
    add_bullets(body.text_frame, bullets, 17)
    if note:
        nbox = slide.shapes.add_textbox(Inches(0.55), Inches(6.2), Inches(8.9), Inches(0.8))
        nbox.text_frame.text = note
        nbox.text_frame.paragraphs[0].font.size = Pt(12)
        nbox.text_frame.paragraphs[0].font.italic = True
        nbox.text_frame.paragraphs[0].font.color.rgb = ACCENT


def add_table_slide(prs, title: str, headers: list[str], rows: list[list[str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.9))
    set_title(tbox, title, 24)
    nrows, ncols = len(rows) + 1, len(headers)
    left, top, width, height = Inches(1.2), Inches(1.6), Inches(7.6), Inches(0.45 * nrows)
    table = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- PART 1: ARC ---
    add_slide_title_only(
        prs,
        "Sau Paper 1: từ measurement gap\nđến decision consequence",
        "Research direction · Paper 2 & 3 (hypothesis-driven)",
    )

    add_content_slide(
        prs,
        "Paper 1 — recap (30 giây)",
        [
            "Benchmark thường: một episode → completion / scalar score",
            "Một episode không trực tiếp test: nếu determining state D đổi, agent có reflect state mới không?",
            "Protocol: giữ instruction, interface, judge; chỉ perturb D (paired G₀, G₁)",
            "Tách ba property: completion · state tracking · score sensitivity / attachment",
            "Ba thứ không đồng nhất (Type A / score-sensitive / high-score miss)",
        ],
    )

    add_content_slide(
        prs,
        "Kết luận Paper 1 — claim hẹp",
        [
            "Không phải: “Benchmark này broken.”",
            "Mà: High score không certify agent đã grounded vào determining state hiện tại",
            "Completion, tracking, score attachment = các property có thể tách biệt",
        ],
        note="Phenomenon discovery — chưa trả lời “for what?”",
    )

    add_content_slide(
        prs,
        "Câu hỏi tiếp theo",
        [
            "Không phải: “Em sẽ đề xuất metric mới.”",
            "Mà: Measurement gap có ảnh hưởng quyết định chọn agent để deploy không?",
            "Benchmark được dùng để: so sánh · chọn model · claim capability · quyết deploy",
            "→ Test decision consequence trước khi xây framework",
        ],
    )

    add_content_slide(
        prs,
        "Research arc (3 bước)",
        [
            "Paper 1 — Discovery: Is there a measurement gap?",
            "Paper 2 — Consequence: Does the gap affect which agent we choose?",
            "Paper 3 — Solution: How should we evaluate? (chỉ nếu Paper 2 đứng)",
            "",
            "discover → establish consequence → build infrastructure",
        ],
    )

    add_content_slide(
        prs,
        "Paper 2 — preview (decision)",
        [
            "Quyết định: chọn 1 CUA cho môi trường stateful / state thay đổi",
            "Rule 1 — Winner_Score = argmax average benchmark score trên G₀",
            "Rule 2 — Winner_STS = argmax state tracking trên paired interventions",
            "Primary question: Hai winner có cùng một agent không?",
        ],
    )

    add_content_slide(
        prs,
        "Paper 3 — preview (conditional)",
        [
            "Chỉ justified nếu Paper 2 có selection consequence",
            "Primitive: State-grounded reliability không establish từ một episode",
            "Cần controlled perturbation của determining state D",
            "Protocol layer cho lab / benchmark maintainer — không rank mọi CUA",
            "Analog: clean accuracy vs robustness (perturb input)",
        ],
    )

    add_content_slide(
        prs,
        "Arc — one-liner",
        [
            "Paper 1: Does a score tell us what the agent tracked?",
            "Paper 2: Does getting that wrong change which agent we choose?",
            "Paper 3: How should we evaluate once one episode is insufficient?",
        ],
        note="→ Part 2: Paper 2 như decision experiment",
    )

    # --- PART 2: PAPER 2 ---
    add_content_slide(
        prs,
        "Paper 2 — central question",
        [
            "Does benchmark success identify the agent that tracks changing state best?",
            "",
            "Phản biện cần test:",
            "“Score imperfect, nhưng agent score cao vẫn là agent đáng deploy.”",
            "",
            "Operational: argmax S̄⁰  ?=  argmax STS  (cùng analysis set)",
        ],
    )

    add_table_slide(
        prs,
        "Paper 1 vs Paper 2",
        ["", "Paper 1", "Paper 2"],
        [
            ["Mục tiêu", "Tách completion / tracking / ΔS", "Hai decision rule có chọn cùng agent?"],
            ["Outcome", "Type A, sensitive, Type B", "Top-1 selection disagreement"],
            ["Loại study", "Benchmark audit", "Decision experiment"],
        ],
    )

    add_table_slide(
        prs,
        "Ví dụ trực giác (hypothetical — chưa phải kết quả)",
        ["Agent", "Benchmark score", "STS"],
        [
            ["A", "92", "70"],
            ["B", "89", "90"],
            ["C", "85", "78"],
        ],
    )

    add_content_slide(
        prs,
        "Không giả định inversion trước",
        [
            "Paper 1: exploratory signal — không confirmatory ranking",
            "Paper 2: fresh pre-specified task universe T_confirmatory",
            "Paper 1’s 10 tasks: replay pilot Layer A only — không vào confirmatory rank",
            "Agent set M: sealed IDs trước khi chạy — không drop sau outcome",
        ],
    )

    add_content_slide(
        prs,
        "Layer A — Calibration / identifiability",
        [
            "Score cao có predict P(track = 1) cao hơn không?",
            "Across paired cells: S → tracking",
            "Null mạnh nếu align tốt:",
            "  → Gap tồn tại ở episode level nhưng aggregate score vẫn informative",
            "Correlation thấp ≠ contribution cuối (chưa đổi top-1)",
        ],
    )

    add_content_slide(
        prs,
        "Layer B — Selection (primary)",
        [
            "M_score = argmax_M  S̄⁰(M)   — base leg only, như leaderboard",
            "M_STS  = argmax_M  STS(M)  — cùng valid-pair analysis set",
            "Primary outcome: M_score  ?=  M_STS",
            "Không average S⁰ và S¹ vào “success score”",
            "n_min = 3 valid pairs để vào argmax; zero valid = execution coverage",
        ],
    )

    add_content_slide(
        prs,
        "Các scenario (đều publishable)",
        [
            "1. Align: top-1 giống → gap chưa có decision harm rõ",
            "2. Calibration yếu, top-1 giống → measurement ≠ decision (middle)",
            "3. Top-1 khác → có thể chọn sai agent (strongest)",
            "4. (exploratory) Rank khác theo state family → cần profile, không một scalar",
        ],
    )

    add_content_slide(
        prs,
        "Đơn vị thí nghiệm (M, T, I)",
        [
            "G₁ = I(G₀): chỉ determining set D đổi",
            "Giữ: instruction · UI · task structure · judge",
            "Đo: completion · STS (typed D, guest gold) · S⁰, S¹",
            "ΔS = score attachment audit — không phải reliability metric",
        ],
    )

    add_content_slide(
        prs,
        "Cross-agent & cross-task",
        [
            "Agents: ≥4, sealed trước run; không publish subset invert",
            "Tasks: universe T mới — eligibility rule trong PAPER2_SPEC",
            "Stratify state families: numeric · categorical · aggregation · temporal · joint · entity",
            "N suy ra từ frozen rule — không KPI “30 task × 5 model”",
            "Primary: 1 CF/task; ~25% subset thêm G₂ (robustness check)",
        ],
    )

    add_content_slide(
        prs,
        "Intervention discipline",
        [
            "Paper 1: designed role ≠ observed class",
            "Paper 2: intervention gold-moving + identifiable — outcome do experiment quyết định",
            "Không “trap benchmark” / không design để ép Type B",
        ],
    )

    add_content_slide(
        prs,
        "Frozen trước khi chạy (PAPER2_SPEC.md)",
        [
            "1. Primary hypothesis (top-1 disagreement)",
            "2. Target deployment decision (chọn 1 CUA stateful)",
            "3. Agent inclusion rule",
            "4. Task / state-family inclusion rule",
            "5. STS + S̄⁰ definition",
            "6. Selection disagreement + null scenarios",
        ],
        note="Chưa chạy cell mới — bước tiếp: sealed M và T",
    )

    add_content_slide(
        prs,
        "Primitive (thesis sâu hơn metric)",
        [
            "State-grounded reliability cannot be established from one realized episode.",
            "Need: controlled perturbation of determining state D.",
            "Reporting task success alone may be insufficient for stateful deployment.",
            "STS = quantity sau khi chấp nhận primitive — không phải “metric vì hay”",
        ],
    )

    add_content_slide(
        prs,
        "Câu hỏi cho anh Khải",
        [
            "Arc Paper 1 → 2 → 3 có đủ lớn và đúng “for what?” không?",
            "Paper 2 (top-1 disagreement trên frozen universe) có đủ decision utility không?",
            "Em có cần deployment consequence mạnh hơn trước khi chạy — hay Level 3 để sau?",
            "",
            "Em sẵn sàng nhận null nếu score vẫn đủ chọn agent — nhưng muốn test fair trước.",
        ],
    )

    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
